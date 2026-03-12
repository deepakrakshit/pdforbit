"""
powerpoint_to_pdf.py — Enterprise-grade PowerPoint to PDF Processor
====================================================================
KEY FEATURES:
  • LibreOffice primary path: pixel-perfect fidelity preservation
    (slide visuals, images, charts, animations stripped cleanly)
  • python-pptx fallback: renders each slide as a proportional page with
    title + body structure and optional speaker notes appendix
  • Slide dimensions respected (16:9, 4:3, etc.) in fallback output
  • Speaker notes appendix when notes are present
  • Conversion engine reported in metadata
"""
from __future__ import annotations

import logging
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.util import Emu

from app.models.enums import ArtifactKind
from app.schemas.job import ConvertToPdfJobRequest
from app.services.pdf.advanced_utils import pdf_page_count
from app.services.pdf.common import (
    BaseToolProcessor,
    GeneratedArtifact,
    PDF_CONTENT_TYPE,
    PdfProcessingError,
    ProcessingResult,
    ProcessorContext,
    ensure_pdf_output_filename,
)
from app.utils.libreoffice import LibreOfficeConversionError, LibreOfficeUnavailableError, convert_with_libreoffice

log = logging.getLogger(__name__)

_POINTS_PER_EMU = 1.0 / 914400.0 * 72.0
_MARGIN = 36.0
_FONT_TITLE = 20.0
_FONT_BODY = 11.0
_LINE_H_TITLE = 28.0
_LINE_H_BODY = 16.0


class PowerPointToPdfProcessor(BaseToolProcessor):
    tool_id = "ppt2pdf"

    def process(self, context: ProcessorContext) -> ProcessingResult:
        payload = ConvertToPdfJobRequest.model_validate(context.payload)
        source = context.require_single_input()
        output_filename = ensure_pdf_output_filename(payload.output_filename)
        output_path = context.workspace / output_filename
        include_notes = getattr(payload, "include_speaker_notes", False)

        # --- Primary: LibreOffice ---
        try:
            converted_path = convert_with_libreoffice(
                source.storage_path,
                output_dir=context.workspace,
                target_format="pdf",
            )
            if converted_path != output_path:
                converted_path.replace(output_path)
            return ProcessingResult(
                artifact=GeneratedArtifact(
                    local_path=output_path,
                    filename=output_filename,
                    content_type=PDF_CONTENT_TYPE,
                    kind=ArtifactKind.RESULT,
                    metadata={
                        "pages_processed": pdf_page_count(output_path),
                        "conversion_engine": "libreoffice",
                    },
                ),
                completion_message="Presentation converted to PDF successfully.",
            )
        except LibreOfficeUnavailableError:
            log.info("ppt2pdf: LibreOffice not available, using python-pptx fallback")
        except LibreOfficeConversionError as exc:
            raise PdfProcessingError(
                code="presentation_conversion_failed",
                user_message=f"Could not convert the presentation: {exc}",
            ) from exc

        # --- Fallback: python-pptx structured renderer ---
        try:
            prs = Presentation(source.storage_path)
        except Exception as exc:
            raise PdfProcessingError(
                code="ppt_read_failed",
                user_message="Could not open the presentation file. It may be corrupted.",
            ) from exc

        # Determine slide dimensions in points
        slide_w_pt = float(prs.slide_width) * _POINTS_PER_EMU if prs.slide_width else 720.0
        slide_h_pt = float(prs.slide_height) * _POINTS_PER_EMU if prs.slide_height else 540.0

        pdf_doc = fitz.open()
        notes_pages: list[tuple[int, str]] = []  # (slide_number, notes_text)

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            _render_slide(slide, slide_num, pdf_doc, slide_w_pt=slide_w_pt, slide_h_pt=slide_h_pt)

            # Collect speaker notes
            if include_notes and slide.has_notes_slide:
                notes_text = _extract_notes(slide)
                if notes_text:
                    notes_pages.append((slide_num, notes_text))

        # Append speaker notes as an appendix
        if notes_pages:
            _render_notes_appendix(notes_pages, pdf_doc, page_w=slide_w_pt, page_h=slide_h_pt)

        if pdf_doc.page_count == 0:
            pdf_doc.new_page(width=720.0, height=540.0)

        pdf_doc.save(output_path, garbage=4, deflate=True)
        pdf_doc.close()

        return ProcessingResult(
            artifact=GeneratedArtifact(
                local_path=output_path,
                filename=output_filename,
                content_type=PDF_CONTENT_TYPE,
                kind=ArtifactKind.RESULT,
                metadata={
                    "pages_processed": pdf_page_count(output_path),
                    "conversion_engine": "python-pptx-fallback",
                    "notes_included": include_notes and bool(notes_pages),
                },
            ),
            completion_message="Presentation converted to PDF successfully.",
        )


# ---------------------------------------------------------------------------
# Slide renderer
# ---------------------------------------------------------------------------

def _render_slide(slide, slide_num: int, pdf: fitz.Document, *, slide_w_pt: float, slide_h_pt: float) -> None:
    """Renders a single slide as a PDF page with title and body text."""
    page = pdf.new_page(width=slide_w_pt, height=slide_h_pt)

    # Light grey background (better than pure white for presentation slides)
    page.draw_rect(fitz.Rect(0, 0, slide_w_pt, slide_h_pt), color=None, fill=(0.97, 0.97, 0.97))

    # Slide number badge in top-right corner
    badge_text = str(slide_num)
    page.insert_text(
        fitz.Point(slide_w_pt - _MARGIN - 20, _MARGIN - 5),
        badge_text,
        fontname="helv",
        fontsize=8.0,
        color=(0.6, 0.6, 0.6),
    )

    y = _MARGIN
    title_done = False

    title_shape = getattr(slide.shapes, "title", None)
    title_text = ""
    if title_shape is not None:
        title_text = (getattr(title_shape, "text", "") or "").strip()

    if title_text:
        title_done = True
        page.draw_rect(
            fitz.Rect(_MARGIN - 4, y - 2, slide_w_pt - _MARGIN + 4, y + _LINE_H_TITLE + 2),
            color=None,
            fill=(0.15, 0.35, 0.65),
        )
        page.insert_text(
            fitz.Point(_MARGIN, y + (_FONT_TITLE * 0.95)),
            title_text,
            fontname="helv",
            fontsize=_FONT_TITLE,
            color=(1.0, 1.0, 1.0),
        )
        y += _LINE_H_TITLE + 10

    for shape in slide.shapes:
        if title_shape is not None and shape == title_shape:
            continue
        if not hasattr(shape, "text_frame"):
            continue
        text = shape.text.strip()
        if not text:
            continue

        is_title = (not title_done and shape.shape_type in (13, 14)) or (
            not title_done and hasattr(shape, "name") and "title" in (shape.name or "").lower()
        )

        if is_title:
            title_done = True
            # Draw title box
            page.draw_rect(
                fitz.Rect(_MARGIN - 4, y - 2, slide_w_pt - _MARGIN + 4, y + _LINE_H_TITLE + 2),
                color=None,
                fill=(0.15, 0.35, 0.65),
            )
            page.insert_textbox(
                fitz.Rect(_MARGIN, y, slide_w_pt - _MARGIN, y + _LINE_H_TITLE),
                text,
                fontname="helv",
                fontsize=_FONT_TITLE,
                color=(1.0, 1.0, 1.0),
                align=0,
            )
            y += _LINE_H_TITLE + 10
        else:
            # Body text — handle bullet points from text frame paragraphs
            for para in shape.text_frame.paragraphs:
                para_text = para.text.strip()
                if not para_text:
                    y += 4
                    continue
                level = para.level or 0
                indent = _MARGIN + level * 12
                bullet = "• " if level == 0 else "  – "
                display = f"{bullet}{para_text}"
                result = page.insert_textbox(
                    fitz.Rect(indent, y, slide_w_pt - _MARGIN, y + _LINE_H_BODY * 2),
                    display,
                    fontname="helv",
                    fontsize=_FONT_BODY - level,
                    color=(0.1, 0.1, 0.1),
                    align=0,
                )
                y += _LINE_H_BODY
                if y > slide_h_pt - _MARGIN:
                    break


def _extract_notes(slide) -> str:
    """Extracts speaker notes text from a slide."""
    try:
        notes_tf = slide.notes_slide.notes_text_frame
        return notes_tf.text.strip()
    except Exception:
        return ""


def _render_notes_appendix(
    notes_pages: list[tuple[int, str]],
    pdf: fitz.Document,
    *,
    page_w: float,
    page_h: float,
) -> None:
    """Appends a speaker notes section at the end of the PDF."""
    page = pdf.new_page(width=page_w, height=page_h)
    y = _MARGIN
    page.insert_text(fitz.Point(_MARGIN, y + 16), "Speaker Notes", fontname="helv", fontsize=16.0, color=(0, 0, 0))
    y += 28

    for slide_num, notes_text in notes_pages:
        if y + _LINE_H_BODY * 3 > page_h - _MARGIN:
            page = pdf.new_page(width=page_w, height=page_h)
            y = _MARGIN

        page.insert_text(fitz.Point(_MARGIN, y + 11), f"Slide {slide_num}:", fontname="helv", fontsize=10.0, color=(0.2, 0.2, 0.8))
        y += 14
        result = page.insert_textbox(
            fitz.Rect(_MARGIN, y, page_w - _MARGIN, page_h - _MARGIN),
            notes_text,
            fontname="helv",
            fontsize=9.0,
            color=(0.2, 0.2, 0.2),
        )
        y += max(0, -result) + 8 if result < 0 else _LINE_H_BODY * max(1, notes_text.count("\n") + 1) + 8
        y += 6