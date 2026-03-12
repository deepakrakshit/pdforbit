from __future__ import annotations

import io
from pathlib import Path

from docx import Document as WordDocument
import fitz
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from fastapi.testclient import TestClient
import pytest

from app.utils.libreoffice import LibreOfficeUnavailableError

from tests.support.integration import (
    authenticate_user,
    create_migrated_client,
    get_primary_artifact,
    run_queued_worker,
)


@pytest.fixture()
def pdf_convert_to_pdf_client(tmp_path: Path, backend_root: Path) -> TestClient:
    client = create_migrated_client(
        tmp_path=tmp_path,
        backend_root=backend_root,
        database_name="pdf-convert-to-pdf.sqlite3",
        storage_name="storage",
        access_secret="pdf-convert-to-pdf-access-secret-with-32-plus-chars",
        refresh_secret="pdf-convert-to-pdf-refresh-secret-with-32-plus-chars",
        settings_overrides={"TRANSLATION_PROVIDER": "mock"},
    )
    client.headers.update(authenticate_user(client, email="pdf-convert-to-pdf@pdforbit.test"))
    return client


def test_ocr_job_generates_searchable_pdf(
    pdf_convert_to_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    file_id = upload_file(
        pdf_convert_to_pdf_client,
        filename="scan.pdf",
        content=image_pdf_bytes("HELLO OCR"),
        content_type="application/pdf",
    )

    job_id, poll_body = submit_and_run_job(
        pdf_convert_to_pdf_client,
        monkeypatch,
        endpoint="/api/v1/optimize/ocr",
        payload={
            "file_id": file_id,
            "language": "eng",
            "dpi": 300,
            "output_filename": "searchable.pdf",
        },
    )

    assert poll_body["pages_processed"] == 1
    _, artifact_path = get_primary_artifact(pdf_convert_to_pdf_client, job_id=job_id)
    with fitz.open(artifact_path) as output_pdf:
        assert output_pdf.page_count == 1
        assert len(output_pdf[0].get_text("text").strip()) > 0


def test_img2pdf_job_generates_pdf(pdf_convert_to_pdf_client: TestClient, monkeypatch: pytest.MonkeyPatch) -> None:
    file_id = upload_file(
        pdf_convert_to_pdf_client,
        filename="image.png",
        content=png_bytes("Image to PDF"),
        content_type="image/png",
    )

    job_id, _ = submit_and_run_job(
        pdf_convert_to_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/to-pdf",
        payload={"file_id": file_id, "output_filename": "image.pdf"},
    )

    _, artifact_path = get_primary_artifact(pdf_convert_to_pdf_client, job_id=job_id)
    with fitz.open(artifact_path) as output_pdf:
        assert output_pdf.page_count == 1


def test_img2pdf_job_supports_multiple_images(pdf_convert_to_pdf_client: TestClient, monkeypatch: pytest.MonkeyPatch) -> None:
    first_file_id = upload_file(
        pdf_convert_to_pdf_client,
        filename="image-a.png",
        content=png_bytes("Image A"),
        content_type="image/png",
    )
    second_file_id = upload_file(
        pdf_convert_to_pdf_client,
        filename="image-b.png",
        content=png_bytes("Image B"),
        content_type="image/png",
    )

    job_id, poll_body = submit_and_run_job(
        pdf_convert_to_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/to-pdf",
        payload={
            "file_ids": [first_file_id, second_file_id],
            "dpi": 300,
            "page_size": "Letter",
            "output_filename": "album.pdf",
        },
    )

    assert poll_body["pages_processed"] == 2
    artifact, artifact_path = get_primary_artifact(pdf_convert_to_pdf_client, job_id=job_id)
    assert artifact.metadata_json["page_size"] == "Letter"
    assert artifact.metadata_json["dpi"] == 300
    with fitz.open(artifact_path) as output_pdf:
        assert output_pdf.page_count == 2


def test_word2pdf_job_generates_pdf(pdf_convert_to_pdf_client: TestClient, monkeypatch: pytest.MonkeyPatch) -> None:
    file_id = upload_file(
        pdf_convert_to_pdf_client,
        filename="document.docx",
        content=docx_bytes(["Quarterly report", "Revenue increased"]),
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )

    job_id, _ = submit_and_run_job(
        pdf_convert_to_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/to-pdf",
        payload={"file_id": file_id, "output_filename": "document.pdf"},
    )

    _, artifact_path = get_primary_artifact(pdf_convert_to_pdf_client, job_id=job_id)
    with fitz.open(artifact_path) as output_pdf:
        assert "Quarterly report" in output_pdf[0].get_text("text")


def test_excel2pdf_job_generates_pdf(pdf_convert_to_pdf_client: TestClient, monkeypatch: pytest.MonkeyPatch) -> None:
    file_id = upload_file(
        pdf_convert_to_pdf_client,
        filename="sheet.xlsx",
        content=xlsx_bytes([["Region", "Revenue"], ["North", 1200], ["South", 980]]),
        content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )

    job_id, _ = submit_and_run_job(
        pdf_convert_to_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/to-pdf",
        payload={"file_id": file_id, "output_filename": "sheet.pdf"},
    )

    _, artifact_path = get_primary_artifact(pdf_convert_to_pdf_client, job_id=job_id)
    with fitz.open(artifact_path) as output_pdf:
        text = output_pdf[0].get_text("text")
        assert "Region" in text
        assert "Revenue" in text


def test_ppt2pdf_job_generates_pdf(pdf_convert_to_pdf_client: TestClient, monkeypatch: pytest.MonkeyPatch) -> None:
    file_id = upload_file(
        pdf_convert_to_pdf_client,
        filename="slides.pptx",
        content=pptx_bytes(["Slide One", "Slide Two"]),
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    job_id, _ = submit_and_run_job(
        pdf_convert_to_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/to-pdf",
        payload={"file_id": file_id, "output_filename": "slides.pdf"},
    )

    _, artifact_path = get_primary_artifact(pdf_convert_to_pdf_client, job_id=job_id)
    with fitz.open(artifact_path) as output_pdf:
        assert output_pdf.page_count >= 1
        assert "Slide One" in output_pdf[0].get_text("text")


def test_ppt2pdf_job_can_include_speaker_notes_in_fallback(
    pdf_convert_to_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import app.services.pdf.powerpoint_to_pdf as ppt_to_pdf_module

    monkeypatch.setattr(
        ppt_to_pdf_module,
        "convert_with_libreoffice",
        lambda *args, **kwargs: (_ for _ in ()).throw(LibreOfficeUnavailableError("missing")),
    )

    file_id = upload_file(
        pdf_convert_to_pdf_client,
        filename="notes.pptx",
        content=pptx_bytes([("Strategy", "Discuss runway and margin expansion")]),
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    job_id, _ = submit_and_run_job(
        pdf_convert_to_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/to-pdf",
        payload={
            "file_id": file_id,
            "include_speaker_notes": True,
            "output_filename": "notes.pdf",
        },
    )

    artifact, artifact_path = get_primary_artifact(pdf_convert_to_pdf_client, job_id=job_id)
    assert artifact.metadata_json["conversion_engine"] == "python-pptx-fallback"
    assert artifact.metadata_json["notes_included"] is True
    with fitz.open(artifact_path) as output_pdf:
        full_text = "\n".join(page.get_text("text") for page in output_pdf)
        assert "Speaker Notes" in full_text
        assert "margin expansion" in full_text


def test_html2pdf_job_generates_pdf(pdf_convert_to_pdf_client: TestClient, monkeypatch: pytest.MonkeyPatch) -> None:
    file_id = upload_file(
        pdf_convert_to_pdf_client,
        filename="page.html",
        content=b"<html><body><h1>Hello</h1><p>HTML to PDF</p></body></html>",
        content_type="text/html",
    )

    job_id, _ = submit_and_run_job(
        pdf_convert_to_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/html-to-pdf",
        payload={"file_id": file_id, "page_size": "A4", "output_filename": "page.pdf"},
    )

    _, artifact_path = get_primary_artifact(pdf_convert_to_pdf_client, job_id=job_id)
    with fitz.open(artifact_path) as output_pdf:
        text = output_pdf[0].get_text("text")
        assert "Hello" in text
        assert "HTML to PDF" in text


def test_html2pdf_job_supports_url_source(pdf_convert_to_pdf_client: TestClient, monkeypatch: pytest.MonkeyPatch) -> None:
    import app.services.pdf.html_to_pdf as html_to_pdf_module

    def fake_download(url: str, *, workspace: Path) -> Path:
        html_path = workspace / "downloaded.html"
        html_path.write_text("<html><body><h1>Fetched</h1><p>Remote content</p></body></html>", encoding="utf-8")
        return html_path

    monkeypatch.setattr(html_to_pdf_module.HtmlToPdfProcessor, "_download_url_html", staticmethod(fake_download))

    job_id, _ = submit_and_run_job(
        pdf_convert_to_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/html-to-pdf",
        payload={
            "url": "https://example.com/report",
            "page_size": "A4",
            "output_filename": "fetched.pdf",
        },
    )

    artifact, artifact_path = get_primary_artifact(pdf_convert_to_pdf_client, job_id=job_id)
    assert artifact.metadata_json["source_url"] == "https://example.com/report"
    with fitz.open(artifact_path) as output_pdf:
        text = output_pdf[0].get_text("text")
        assert "Fetched" in text
        assert "Remote content" in text


def submit_and_run_job(
    client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
    *,
    endpoint: str,
    payload: dict[str, object],
) -> tuple[str, dict[str, object]]:
    response = client.post(endpoint, json=payload)
    assert response.status_code == 201
    job_id = response.json()["job_id"]
    run_queued_worker(client, monkeypatch)
    poll_response = client.get(f"/api/v1/jobs/{job_id}")
    assert poll_response.status_code == 200
    poll_body = poll_response.json()
    assert poll_body["status"] == "completed"
    assert poll_body["progress"] == 100
    return job_id, poll_body


def upload_file(client: TestClient, *, filename: str, content: bytes, content_type: str) -> str:
    response = client.post("/api/v1/upload", files={"file": (filename, content, content_type)})
    assert response.status_code == 201
    return response.json()["file_id"]


def png_bytes(text: str) -> bytes:
    image = Image.new("RGB", (900, 280), "white")
    draw = ImageDraw.Draw(image)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 72)
    except OSError:
        font = ImageFont.load_default()
    draw.text((40, 90), text, fill="black", font=font)
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def image_pdf_bytes(text: str) -> bytes:
    image = Image.open(io.BytesIO(png_bytes(text)))
    buffer = io.BytesIO()
    image.convert("RGB").save(buffer, format="PDF", resolution=150.0)
    return buffer.getvalue()


def docx_bytes(paragraphs: list[str]) -> bytes:
    document = WordDocument()
    for paragraph in paragraphs:
        document.add_paragraph(paragraph)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def xlsx_bytes(rows: list[list[object]]) -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    for row in rows:
        sheet.append(row)
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def pptx_bytes(slides: list[str] | list[tuple[str, str]]) -> bytes:
    presentation = Presentation()
    for slide_spec in slides:
        if isinstance(slide_spec, tuple):
            slide_text, notes_text = slide_spec
        else:
            slide_text, notes_text = slide_spec, ""
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = slide_text
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()