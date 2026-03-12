from __future__ import annotations

import zipfile
from pathlib import Path

from docx import Document as WordDocument
import fitz
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from fastapi.testclient import TestClient
import pytest

from tests.support.integration import (
    authenticate_user,
    create_migrated_client,
    get_primary_artifact,
    run_queued_worker,
    upload_pdf,
)


@pytest.fixture()
def pdf_convert_from_pdf_client(tmp_path: Path, backend_root: Path) -> TestClient:
    client = create_migrated_client(
        tmp_path=tmp_path,
        backend_root=backend_root,
        database_name="pdf-convert-from-pdf.sqlite3",
        storage_name="storage",
        access_secret="pdf-convert-from-pdf-access-secret-with-32-plus-chars",
        refresh_secret="pdf-convert-from-pdf-refresh-secret-with-32-plus-chars",
        settings_overrides={"TRANSLATION_PROVIDER": "mock"},
    )
    client.headers.update(authenticate_user(client, email="pdf-convert-from-pdf@pdforbit.test"))
    return client


def test_compare_job_generates_diff_archive(
    pdf_convert_from_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    first_file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["Same text"]))
    second_file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["Changed text"]))

    job_id, poll_body = submit_and_run_job(
        pdf_convert_from_pdf_client,
        monkeypatch,
        endpoint="/api/v1/security/compare",
        payload={
            "file_id_a": first_file_id,
            "file_id_b": second_file_id,
            "output_filename": "comparison.zip",
        },
    )

    assert poll_body["different_pages"] == 1
    _, artifact_path = get_primary_artifact(pdf_convert_from_pdf_client, job_id=job_id)
    with zipfile.ZipFile(artifact_path) as archive:
        summary = archive.read("summary.txt").decode("utf-8")
        assert "PdfORBIT Comparison Report" in summary
        assert "Pages with differences: 1" in summary
        assert "Changed:" in summary or "Added:" in summary or "Removed:" in summary


def test_pdf2img_job_generates_zip_archive(
    pdf_convert_from_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["One", "Two"]))

    job_id, poll_body = submit_and_run_job(
        pdf_convert_from_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/from-pdf",
        payload={"file_id": file_id, "format": "png", "dpi": 150},
    )

    assert poll_body["parts_count"] == 2
    _, artifact_path = get_primary_artifact(pdf_convert_from_pdf_client, job_id=job_id)
    with zipfile.ZipFile(artifact_path) as archive:
        names = sorted(archive.namelist())
        assert names == ["page-0001.png", "page-0002.png"]


def test_pdf2img_job_supports_single_page_thumbnail(
    pdf_convert_from_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["First", "Second page"]))

    job_id, poll_body = submit_and_run_job(
        pdf_convert_from_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/from-pdf",
        payload={
            "file_id": file_id,
            "format": "webp",
            "dpi": 300,
            "single_page": 2,
            "thumbnail": True,
            "thumbnail_max_px": 180,
        },
    )

    assert poll_body["pages_processed"] == 1
    artifact, artifact_path = get_primary_artifact(pdf_convert_from_pdf_client, job_id=job_id)
    assert artifact.content_type == "image/webp"
    with Image.open(artifact_path) as output_image:
        assert max(output_image.size) <= 180


def test_pdf2word_job_generates_docx(
    pdf_convert_from_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["Export to Word"]))

    job_id, _ = submit_and_run_job(
        pdf_convert_from_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/from-pdf",
        payload={"file_id": file_id, "format": "word"},
    )

    _, artifact_path = get_primary_artifact(pdf_convert_from_pdf_client, job_id=job_id)
    document = WordDocument(artifact_path)
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    assert "Export to Word" in text


def test_pdf2word_job_honours_custom_output_filename(
    pdf_convert_from_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["Board minutes"] ))

    job_id, _ = submit_and_run_job(
        pdf_convert_from_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/from-pdf",
        payload={"file_id": file_id, "format": "word", "output_filename": "board-minutes.docx"},
    )

    artifact, artifact_path = get_primary_artifact(pdf_convert_from_pdf_client, job_id=job_id)
    assert artifact.filename == "board-minutes.docx"
    assert artifact_path.name == "board-minutes.docx"


def test_pdf2excel_job_generates_xlsx(
    pdf_convert_from_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["Revenue row"]))

    job_id, _ = submit_and_run_job(
        pdf_convert_from_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/from-pdf",
        payload={"file_id": file_id, "format": "excel"},
    )

    _, artifact_path = get_primary_artifact(pdf_convert_from_pdf_client, job_id=job_id)
    workbook = load_workbook(artifact_path)
    assert workbook.active["A1"].value is not None


def test_pdf2excel_job_honours_custom_output_filename(
    pdf_convert_from_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["Pipeline table"]))

    job_id, _ = submit_and_run_job(
        pdf_convert_from_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/from-pdf",
        payload={"file_id": file_id, "format": "excel", "output_filename": "pipeline.xlsx"},
    )

    artifact, artifact_path = get_primary_artifact(pdf_convert_from_pdf_client, job_id=job_id)
    assert artifact.filename == "pipeline.xlsx"
    assert artifact_path.name == "pipeline.xlsx"


def test_pdf2ppt_job_generates_pptx(
    pdf_convert_from_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["Slide export"]))

    job_id, _ = submit_and_run_job(
        pdf_convert_from_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/from-pdf",
        payload={"file_id": file_id, "format": "ppt"},
    )

    _, artifact_path = get_primary_artifact(pdf_convert_from_pdf_client, job_id=job_id)
    presentation = Presentation(artifact_path)
    slide_text = "\n".join(shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text"))
    assert "Slide export" in slide_text


def test_pdf2ppt_job_honours_custom_output_filename(
    pdf_convert_from_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["Quarterly deck"]))

    job_id, _ = submit_and_run_job(
        pdf_convert_from_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/from-pdf",
        payload={"file_id": file_id, "format": "ppt", "output_filename": "quarterly-deck.pptx"},
    )

    artifact, artifact_path = get_primary_artifact(pdf_convert_from_pdf_client, job_id=job_id)
    assert artifact.filename == "quarterly-deck.pptx"
    assert artifact_path.name == "quarterly-deck.pptx"


def test_pdf2pdfa_job_generates_archival_pdf(
    pdf_convert_from_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["Archive me"]))

    job_id, _ = submit_and_run_job(
        pdf_convert_from_pdf_client,
        monkeypatch,
        endpoint="/api/v1/convert/from-pdf",
        payload={"file_id": file_id, "pdfa_level": "2b", "output_filename": "archive.pdf"},
    )

    artifact, artifact_path = get_primary_artifact(pdf_convert_from_pdf_client, job_id=job_id)
    assert artifact.metadata_json["pdfa_level"] == "2b"
    with fitz.open(artifact_path) as output_pdf:
        assert output_pdf.page_count == 1


def test_compare_job_supports_text_diff_mode(
    pdf_convert_from_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    first_file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["Alpha Beta"] ))
    second_file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["Alpha Gamma"] ))

    job_id, _ = submit_and_run_job(
        pdf_convert_from_pdf_client,
        monkeypatch,
        endpoint="/api/v1/security/compare",
        payload={
            "file_id_a": first_file_id,
            "file_id_b": second_file_id,
            "diff_mode": "text",
            "output_filename": "text-comparison.zip",
        },
    )

    artifact, artifact_path = get_primary_artifact(pdf_convert_from_pdf_client, job_id=job_id)
    assert artifact.metadata_json["diff_mode"] == "text"
    with zipfile.ZipFile(artifact_path) as archive:
        names = set(archive.namelist())
        assert "document-1-annotated.pdf" in names
        assert "document-2-annotated.pdf" in names
        summary = archive.read("summary.txt").decode("utf-8")
        assert "Changed" in summary or "Added" in summary


def test_translate_job_generates_mock_translated_pdf(
    pdf_convert_from_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    file_id = upload_pdf(pdf_convert_from_pdf_client, pdf_bytes=build_text_pdf_bytes(["Hello translation"]))

    job_id, poll_body = submit_and_run_job(
        pdf_convert_from_pdf_client,
        monkeypatch,
        endpoint="/api/v1/intelligence/translate",
        payload={"file_id": file_id, "target_language": "fr", "source_language": "en"},
    )

    assert poll_body["detected_language"] == "en"
    assert poll_body["word_count"] >= 2
    _, artifact_path = get_primary_artifact(pdf_convert_from_pdf_client, job_id=job_id)
    with fitz.open(artifact_path) as output_pdf:
        assert "[en->fr]" in output_pdf[0].get_text("text")


def test_summarize_job_generates_pdf_brief(
    pdf_convert_from_pdf_client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    file_id = upload_pdf(
        pdf_convert_from_pdf_client,
        pdf_bytes=build_text_pdf_bytes([
            "Quarterly revenue grew 22 percent year over year.",
            "Main risks were churn in the SMB segment and delayed collections.",
        ]),
    )

    job_id, poll_body = submit_and_run_job(
        pdf_convert_from_pdf_client,
        monkeypatch,
        endpoint="/api/v1/intelligence/summarize",
        payload={
            "file_id": file_id,
            "output_language": "en",
            "length": "short",
            "focus": "Highlight growth, risk, and collections",
        },
    )

    assert poll_body["word_count"] >= 10
    _, artifact_path = get_primary_artifact(pdf_convert_from_pdf_client, job_id=job_id)
    with fitz.open(artifact_path) as output_pdf:
        text = output_pdf[0].get_text("text")
        assert "EXECUTIVE SUMMARY" in text
        assert "summary:short:en" in text


def submit_and_run_job(
    client: TestClient,
    monkeypatch: pytest.MonkeyPatch,
    *,
    endpoint: str,
    payload: dict[str, object],
) -> tuple[str, dict[str, object]]:
    response = client.post(endpoint, json=payload)
    assert response.status_code == 201
    job_id = response.json()["job_id"]
    run_queued_worker(client, monkeypatch)
    poll_response = client.get(f"/api/v1/jobs/{job_id}")
    assert poll_response.status_code == 200
    poll_body = poll_response.json()
    assert poll_body["status"] == "completed"
    assert poll_body["progress"] == 100
    return job_id, poll_body


def build_text_pdf_bytes(texts: list[str]) -> bytes:
    document = fitz.open()
    for text in texts:
        page = document.new_page(width=595, height=842)
        page.insert_text((72, 120), text, fontsize=24)
    return document.tobytes()